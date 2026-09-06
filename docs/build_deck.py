#!/usr/bin/env python3
"""雨アラートのしくみ — 説明資料(PowerPoint)を生成する。"""

import sys
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ===== 配色: 気象庁の降水強度カラースケールから =====
NAVY   = C(0x0B, 0x1B, 0x3A)
NAVY2  = C(0x13, 0x29, 0x4F)
BLUE   = C(0x00, 0x41, 0xFF)
BLUE_D = C(0x00, 0x31, 0xC0)
TEAL   = C(0x1C, 0x72, 0x93)
PALE   = C(0xA0, 0xD2, 0xFF)
YELLOW = C(0xFA, 0xF5, 0x00)
RED    = C(0xFF, 0x28, 0x00)
WHITE  = C(0xFF, 0xFF, 0xFF)
LIGHT  = C(0xF4, 0xF7, 0xFD)
INK    = C(0x16, 0x23, 0x3D)
GREY   = C(0x5C, 0x68, 0x84)
LINE   = C(0xD5, 0xDE, 0xEF)
EDGE   = C(0x9F, 0xB0, 0xCE)
ARROW  = C(0x6B, 0x7A, 0x99)
MYFILL = C(0xE4, 0xEC, 0xFF)

FS = "Calibri"
FM = "Courier New"
JP = "Hiragino Sans"   # 日本語用(macOS標準。Windowsでは游ゴシック等に代替される)

prs = Presentation()
prs.slide_width, prs.slide_height = In(13.333), In(7.5)
BLANK = prs.slide_layouts[6]


# ---------- 共通パーツ ----------
def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def text(slide, s, x, y, w, h, *, size=12, bold=False, color=INK, font=FS,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
         spacing=None, char_space=None):
    box = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = Pt(spacing)
        r = p.add_run()
        r.text = line
        f = r.font
        f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
        f.color.rgb = color
        _set_jp_font(r)
        if char_space:  # 字間(python-pptxに無いのでXMLで直接指定)
            r._r.get_or_add_rPr().set("spc", str(int(char_space * 100)))
    return box


def rect(slide, x, y, w, h, *, fill=None, line=None, lw=1.0, dash=False,
         rounded=False, radius=None, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        In(x), In(y), In(w), In(h))
    if rounded and radius is not None:
        shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
        if dash:
            shp.line._get_or_add_ln().append(
                _el('a:prstDash', {'val': 'dash'}))
    if not shadow:
        _no_shadow(shp)
    shp.text_frame.word_wrap = True
    return shp


def ellipse(slide, x, y, w, h, *, fill, line=None, lw=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y), In(w), In(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    _no_shadow(shp)
    return shp


def _el(tag, attrs=None):
    """<a:xxx .../> のような要素を1つ作る。"""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    ns = tag.split(":")[0]
    attr_str = "".join(f' {k}="{v}"' for k, v in (attrs or {}).items())
    return parse_xml(f"<{tag} {nsdecls(ns)}{attr_str}/>")


def _set_jp_font(run):
    """日本語部分のフォントを指定する(未指定だと明朝体になってしまう)。"""
    rPr = run._r.get_or_add_rPr()
    if rPr.find(qn("a:ea")) is None:
        rPr.append(_el("a:ea", {"typeface": JP}))


def _no_shadow(shp):
    """python-pptxの既定の影を消す。"""
    sp = shp._element.spPr
    for tag in ("a:effectLst",):
        for e in sp.findall(qn(tag)):
            sp.remove(e)
    sp.append(_el("a:effectLst"))


def arrow(slide, x1, y1, x2, y2, *, color=ARROW, lw=1.5, both=False, head=True):
    cn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, In(x1), In(y1), In(x2), In(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(lw)
    ln = cn.line._get_or_add_ln()
    if both:
        ln.append(_el("a:headEnd", {"type": "triangle", "w": "med", "len": "med"}))
    if head:
        ln.append(_el("a:tailEnd", {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def dashed_line(slide, x1, y1, x2, y2, color=C(0xC3, 0xCF, 0xE4)):
    cn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, In(x1), In(y1), In(x2), In(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(1)
    cn.line._get_or_add_ln().append(_el("a:prstDash", {"val": "sysDot"}))
    return cn


def card(slide, x, y, w, h):
    return rect(slide, x, y, w, h, fill=WHITE, line=LINE, lw=1.0,
                rounded=True, radius=0.04)


def badge(slide, x, y, label, fill, size=0.46, fsize=13, tcolor=WHITE):
    ellipse(slide, x, y, size, size, fill=fill)
    text(slide, label, x, y, size, size, size=fsize, bold=True,
         color=tcolor, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def light_slide(title, kicker):
    s = prs.slides.add_slide(BLANK)
    bg(s, LIGHT)
    text(s, kicker, 0.62, 0.38, 8, 0.26, size=11, bold=True, color=BLUE, char_space=2)
    text(s, title, 0.6, 0.68, 12.1, 0.62, size=30, bold=True, color=INK)
    return s


def dark_slide():
    s = prs.slides.add_slide(BLANK)
    bg(s, NAVY)
    return s


# =====================================================
# 1. タイトル
# =====================================================
s = dark_slide()
ellipse(s, 9.3, -1.6, 5.6, 5.6, fill=NAVY2)
text(s, "SwitchBot → 気象庁 → LINE", 0.9, 2.05, 8.5, 0.3,
     size=13, bold=True, color=PALE, char_space=3)
text(s, "雨アラートのしくみ", 0.86, 2.45, 9.2, 1.15, size=52, bold=True, color=WHITE)
text(s, "自宅にゲリラ豪雨が近づいたらLINEで知らせ、毎朝6時には\n"
        "その日の雨をグラフで送る。そのしくみと、9つのファイルの役割。",
     0.9, 3.9, 8.8, 0.95, size=15, color=PALE, spacing=26)
for i, t in enumerate(["Python", "GitHub Actions", "LINE Messaging API"]):
    x = 0.9 + i * 2.45
    rect(s, x, 5.2, 2.25, 0.44, fill=NAVY2, line=C(0x2C, 0x44, 0x70),
         rounded=True, radius=0.5)
    text(s, t, x, 5.2, 2.25, 0.44, size=11, color=PALE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================
# 2. できること
# =====================================================
s = light_slide("このシステムがしてくれること", "01 — できること")
for i, it in enumerate([
    ("1", BLUE, "5分ごと", "豪雨をつかまえる",
     "自宅の半径1kmに強い雨が近づくと、降り出す前にLINEへ通知。弱い雨と激しい雨で文面が変わります。"),
    ("2", TEAL, "毎朝 6:00", "1日の雨をグラフで見る",
     "24時間ぶんの降水予報を棒グラフの画像にして送信。傘が要る日かが一目で分かります。"),
    ("3", C(0x2C, 0x5F, 0x7A), "電源オフでもOK", "パソコンが要らない",
     "GitHubのサーバー上で動くため、自分のMacやスマホの電源が入っていなくても動き続けます。"),
]):
    n, col, kick, title, desc = it
    x = 0.6 + i * 4.13
    card(s, x, 1.75, 3.83, 2.85)
    badge(s, x + 0.38, 2.12, n, col)
    text(s, kick, x + 1.0, 2.2, 2.6, 0.3, size=10.5, bold=True, color=col, font=FM)
    text(s, title, x + 0.38, 2.82, 3.2, 0.42, size=17, bold=True, color=INK)
    text(s, desc, x + 0.38, 3.35, 3.1, 1.6, size=12, color=GREY, spacing=20)
text(s, "もともとの出発点は「SwitchBotで何か作れないか」。自宅は動かないので、座標は一度調べれば十分でした。",
     0.62, 5.02, 12, 0.34, size=12, italic=True, color=GREY)

# =====================================================
# 3. システム構成図
# =====================================================
s = light_slide("システム構成図", "02 — 全体像")


def ext_box(x, y, w, h, name, sub):
    rect(s, x, y, w, h, fill=WHITE, line=EDGE, lw=1.5, rounded=True, radius=0.04)
    text(s, name, x, y + 0.15, w, 0.32, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, sub, x, y + 0.5, w, 0.26, size=10, color=GREY, align=PP_ALIGN.CENTER)


def my_box(x, y, w, h, name, sub=None, dash=False):
    rect(s, x, y, w, h, fill=MYFILL, line=BLUE, lw=1.5, dash=dash,
         rounded=True, radius=0.04)
    text(s, name, x, y + (0.14 if sub else 0.19), w, 0.3,
         size=12, bold=True, color=BLUE_D, font=FM, align=PP_ALIGN.CENTER)
    if sub:
        text(s, sub, x, y + 0.48, w, 0.26, size=10, color=GREY, align=PP_ALIGN.CENTER)


def alabel(t, x, y, w):
    text(s, t, x, y, w, 0.24, size=9.5, color=GREY, align=PP_ALIGN.CENTER)


# GitHub Actions の枠(先に置いて背面にする)
rect(s, 3.3, 1.5, 6.3, 4.48, fill=None, line=EDGE, lw=1.25, dash=True,
     rounded=True, radius=0.03)
text(s, "GitHub Actions（クラウド上のサーバー）", 3.45, 1.6, 4.8, 0.24,
     size=10, color=GREY, font=FM)
text(s, "① 5分ごと ─ 豪雨アラート", 3.45, 1.92, 3.4, 0.24, size=10.5, bold=True, color=BLUE)
dashed_line(s, 3.45, 4.60, 9.45, 4.60)
text(s, "② 毎朝6時 ─ 予報グラフ", 3.45, 4.70, 3.4, 0.24, size=10.5, bold=True, color=BLUE)

ext_box(0.5, 2.25, 2.15, 0.9, "気象庁ナウキャスト", "250mメッシュ / 5分更新")
ext_box(0.5, 3.85, 2.15, 0.9, "Open-Meteo", "世界の予報モデル")
my_box(3.50, 2.25, 2.40, 0.9, "jma_nowcast.py", "画像の色 → 雨量")
my_box(6.80, 2.25, 2.40, 0.9, "rain_alert.py", "判定して送信")
my_box(7.30, 3.40, 1.60, 0.62, "state.json", dash=True)
my_box(3.50, 4.98, 2.40, 0.85, "daily_summary.py", "グラフを作る")
ext_box(10.35, 2.20, 2.40, 1.0, "LINE", "スマホに届く")
ext_box(10.35, 4.98, 2.40, 0.85, "imgbb", "画像の置き場所")

arrow(s, 2.65, 2.70, 3.50, 2.70)          # 気象庁 → jma_nowcast
alabel("タイル画像", 2.66, 2.40, 0.83)
arrow(s, 5.90, 2.70, 6.80, 2.70)          # jma_nowcast → rain_alert
alabel("雨量", 5.92, 2.40, 0.86)
arrow(s, 2.65, 4.30, 7.00, 4.30, head=False)  # Open-Meteo → 右へ(途中なので矢印なし)
arrow(s, 7.00, 4.30, 7.00, 3.15)              # 上へ折れて rain_alert へ
alabel("15分ごとの予報", 3.60, 4.02, 1.9)
arrow(s, 8.10, 3.15, 8.10, 3.40, both=True)   # rain_alert ↔ state.json
alabel("通知履歴", 8.95, 3.56, 0.65)
arrow(s, 9.20, 2.70, 10.35, 2.70)         # rain_alert → LINE
alabel("通知を送信", 9.22, 2.40, 1.11)
arrow(s, 5.90, 5.41, 10.35, 5.41)         # daily_summary → imgbb
alabel("グラフ画像を預ける", 7.05, 5.12, 2.0)
arrow(s, 11.55, 4.98, 11.55, 3.20)        # imgbb → LINE
alabel("画像URL", 11.65, 3.95, 1.1)

for i, (fl, ln, t) in enumerate([
        (MYFILL, BLUE, "自分で書いたファイル"),
        (WHITE, EDGE, "外部のサービス")]):
    x = 0.62 + i * 3.1
    rect(s, x, 6.34, 0.3, 0.18, fill=fl, line=ln, lw=1.25, rounded=True, radius=0.15)
    text(s, t, x + 0.42, 6.28, 2.6, 0.3, size=10.5, color=GREY, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================
# 4. 処理をする3ファイル
# =====================================================
s = light_slide("処理をする3つのファイル", "03 — ファイル")
for i, (fn, role, desc) in enumerate([
    ("rain_alert.py", "司令塔",
     "2つのデータ源から雨の強さを取り寄せ、強い方を採用。しきい値（1mm/h・20mm/h）と比べて通知を決め、LINEへ送ります。しきい値もこのファイルの上部にあります。"),
    ("jma_nowcast.py", "気象庁担当",
     "気象庁は数値APIを出しておらず、地図の画像しかありません。緯度経度から画像の位置を計算し、ピクセルの色を読んで雨量に変換します。半径約1kmを調べ最大値を採用。"),
    ("daily_summary.py", "毎朝担当",
     "24時間ぶんの予報から棒グラフの画像を作り、LINEへ送ります。雨の強さで棒の色が変わり、しきい値の位置には点線が入ります。"),
]):
    y = 1.75 + i * 1.55
    card(s, 0.6, y, 12.13, 1.36)
    text(s, fn, 0.95, y + 0.24, 3.2, 0.3, size=14, bold=True, color=BLUE_D, font=FM)
    text(s, role, 0.95, y + 0.64, 3.0, 0.28, size=11, bold=True, color=GREY)
    text(s, desc, 4.3, y + 0.22, 8.05, 1.0, size=12.5, color=INK, spacing=21)

# =====================================================
# 5. 記憶と設定のファイル
# =====================================================
s = light_slide("覚えておく／設定するファイル", "03 — ファイル（続き）")
for i, (fn, role, desc) in enumerate([
    ("rain_alert_state.json", "記憶",
     "前回いつ通知したかだけを書いたメモ。これが無いと同じ雨で何十回も通知が飛びます。"),
    ("requirements.txt", "部品表",
     "使う部品の一覧。requests・Pillow・matplotlib の3つ。実行前にこれを見てそろえます。"),
    ("rain_alert.yml", "時間割",
     "「5分ごとに rain_alert.py を動かして」という、GitHubへの指示書です。"),
    ("daily_summary.yml", "時間割",
     "「毎朝6時に動かして」という指示書。日本語フォントを入れる手順も含みます。"),
    (".env", "秘密（自分のMacのみ）",
     "トークンと自宅の緯度経度。GitHubには絶対に上げず、クラウドではSecretsに登録します。"),
    (".gitignore", "除外リスト",
     "GitHubに上げないファイルの指定。.env を書いておくことで公開事故を防ぎます。"),
]):
    col, row = i % 3, i // 3
    x, y = 0.6 + col * 4.13, 1.75 + row * 2.16
    card(s, x, y, 3.83, 1.92)
    text(s, fn, x + 0.32, y + 0.26, 3.35, 0.3, size=11.5, bold=True, color=BLUE_D, font=FM)
    text(s, role, x + 0.32, y + 0.62, 3.3, 0.26, size=10, bold=True, color=GREY)
    text(s, desc, x + 0.32, y + 0.96, 3.25, 0.85, size=11.5, color=INK, spacing=18)

# =====================================================
# 6. 処理の流れ
# =====================================================
s = light_slide("動くときの流れ", "04 — 順番")
for i, (kick, title, col, steps) in enumerate([
    ("① 5分ごと", "豪雨アラート", BLUE, [
        "GitHubが時間に気づき rain_alert.py を起動",
        "jma_nowcast.py が気象庁の画像から雨量を読む",
        "同時に Open-Meteo からも予報を取り寄せる",
        "2つのうち強い方を採用（見逃しを減らす）",
        "20mm/h以上なら豪雨、1mm/h以上なら雨と判定",
        "前回から時間が空いていればLINEへ送信",
    ]),
    ("② 毎朝 6:00", "予報グラフ", TEAL, [
        "GitHubが daily_summary.py を起動",
        "その日24時間ぶんの降水予報を取得",
        "雨の強さで色分けした棒グラフを作る",
        "画像を imgbb に預けてURLをもらう",
        "URLと一言をまとめてLINEへ送信",
    ]),
]):
    x = 0.6 + i * 6.28
    card(s, x, 1.75, 5.85, 4.35)
    text(s, kick, x + 0.4, 2.02, 3.0, 0.26, size=10.5, bold=True, color=col, font=FM)
    text(s, title, x + 0.4, 2.32, 4.0, 0.4, size=19, bold=True, color=INK)
    for j, st in enumerate(steps):
        y = 2.95 + j * 0.53
        badge(s, x + 0.4, y, str(j + 1), col, size=0.4, fsize=11.5)
        text(s, st, x + 0.98, y, 4.7, 0.4, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================
# 7. 色から雨量を読む
# =====================================================
s = light_slide("色を読んで、雨量に戻す", "05 — 中身")
text(s, "気象庁のナウキャストは、雨の強さを色で表した地図として配信されています。"
        "jma_nowcast.py は自宅の位置にあたるピクセルの色を読み、この対応表で数値に戻しています。",
     0.62, 1.7, 11.9, 0.6, size=13.5, color=GREY, spacing=23)
SCALE = [("F2F2FF", "0.1"), ("A0D2FF", "1"), ("218CFF", "5"), ("0041FF", "10"),
         ("FAF500", "20"), ("FF9900", "30"), ("FF2800", "50"), ("B40068", "80+")]
sw, sx = 1.42, 0.75
for i, (hexc, val) in enumerate(SCALE):
    rect(s, sx + i * sw, 2.62, sw, 1.05,
         fill=C.from_string(hexc), line=C(0xAF, 0xBD, 0xD6), lw=0.75)
    text(s, val, sx + i * sw, 3.75, sw, 0.3, size=12, bold=True,
         color=GREY, font=FM, align=PP_ALIGN.CENTER)
text(s, "mm/h", sx + 8 * sw + 0.08, 3.75, 0.8, 0.3, size=11, color=GREY)
for x0, dotc, msg, tc in [
    (0.75, PALE, "1mm/h 以上 →「まもなく雨」を通知", C(0x1C, 0x5F, 0xA8)),
    (6.9, YELLOW, "20mm/h 以上 →「豪雨注意」を通知", C(0x9A, 0x5B, 0x00)),
]:
    card(s, x0, 4.5, 5.6, 0.85)
    ellipse(s, x0 + 0.32, 4.78, 0.3, 0.3, fill=dotc, line=C(0xAF, 0xBD, 0xD6), lw=0.75)
    text(s, msg, x0 + 0.78, 4.65, 4.7, 0.5, size=13, bold=True, color=tc,
         anchor=MSO_ANCHOR.MIDDLE)
text(s, "この資料の配色も、同じ気象庁のスケールから取っています。",
     0.75, 5.6, 11.8, 0.32, size=12, italic=True, color=GREY)

# =====================================================
# 8. 仕組みのポイント
# =====================================================
s = light_slide("知っておくと分かりやすい4つのこと", "06 — 中身")
for i, (q, h, d) in enumerate([
    ("なぜ2つ？", "データ源を2つ使う理由",
     "気象庁のナウキャストは細かく速い代わりに、公式APIではありません。将来使えなくなる可能性があるため、Open-Meteoを保険として並べています。"),
    ("なぜ間隔をあける？", "同じ雨で連投しない工夫",
     "5分ごとに動くので、雨が1時間続けば12回判定されます。一度送ったら一定時間黙る仕組みを入れました。弱い雨と豪雨は別管理です。"),
    ("実際にあった不具合", "時差という落とし穴",
     "GitHubは世界標準時、天気データは日本時間。そのまま比べて9時間ずれ、雨を全部見逃していました。時刻を日本時間に統一して修正済みです。"),
    ("なぜ画像を外に置く？", "グラフがimgbbを通る理由",
     "LINEには画像そのものではなく「画像のある場所（URL）」を伝える決まりです。作ったグラフはURLを持たないため、imgbbに預けます。"),
]):
    col, row = i % 2, i // 2
    x, y = 0.6 + col * 6.28, 1.75 + row * 2.28
    card(s, x, y, 5.85, 2.04)
    text(s, q, x + 0.38, y + 0.26, 4.8, 0.26, size=10, bold=True, color=GREY, font=FM)
    text(s, h, x + 0.38, y + 0.58, 5.1, 0.36, size=16, bold=True, color=INK)
    text(s, d, x + 0.38, y + 1.02, 5.15, 0.9, size=12, color=GREY, spacing=19)

# =====================================================
# 9. サービスと費用
# =====================================================
s = light_slide("使っているサービスと費用", "07 — 外部サービス")
ROWS = [
    ("サービス", "役割", "費用"),
    ("気象庁ナウキャスト", "250mメッシュ・5分更新の降水データ（画像形式）", "無料・登録不要"),
    ("Open-Meteo", "予報モデルによる降水データ（数値形式）", "無料・登録不要"),
    ("LINE Messaging API", "スマホへの通知", "無料枠 月200通"),
    ("imgbb", "グラフ画像の置き場所", "無料"),
    ("GitHub Actions", "決まった時刻にプログラムを動かすサーバー", "公開リポジトリなら無制限"),
]
tbl = s.shapes.add_table(len(ROWS), 3, In(0.6), In(1.78), In(12.13), In(3.1)).table
for w, cw in zip(tbl.columns, (In(3.1), In(6.03), In(3.0))):
    w.width = cw
for ri, row in enumerate(ROWS):
    tbl.rows[ri].height = In(0.5)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = In(0.14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY2 if ri == 0 else WHITE
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = val
        r.font.size = Pt(12)
        r.font.name = FS
        r.font.bold = (ri == 0) or (ci == 0)
        r.font.color.rgb = WHITE if ri == 0 else (INK if ci == 0 else GREY)
        _set_jp_font(r)
text(s, "GitHub Actionsは1回の実行を1分単位に切り上げて計上します。5分間隔だと月8,640分に達し、"
        "非公開リポジトリの無料枠2,000分では約7日で尽きるため、公開リポジトリに切り替えました。",
     0.62, 5.1, 12, 0.6, size=12, color=GREY, spacing=20)
text(s, "LINE Notify は2025年3月末で終了しているため、後継の Messaging API を使っています。",
     0.62, 5.75, 12, 0.32, size=12, italic=True, color=GREY)

# =====================================================
# 10. 注意点
# =====================================================
s = dark_slide()
ellipse(s, -1.5, 6.45, 4.0, 4.0, fill=NAVY2)
text(s, "08 — 注意", 0.62, 0.5, 8, 0.26, size=11, bold=True, color=PALE, char_space=2)
text(s, "運用するうえで気をつけること", 0.6, 0.8, 11, 0.6, size=30, bold=True, color=WHITE)
for i, (dotc, title, desc) in enumerate([
    (RED, "座標をコードに書き戻さない",
     "リポジトリは公開しています。緯度経度は .env とSecretsにだけ。一度コミットすると、消しても履歴に残り続けます。"),
    (YELLOW, "LINEの無料枠は月200通",
     "毎朝のグラフが月60通を使うため、豪雨アラートに使えるのは月140通ほど。超えるとその月は送れません。"),
    (PALE, "実行間隔は5分（3分ではない）",
     "GitHub Actionsの定期実行は最短5分が目安。気象庁のデータ自体が5分更新なので鮮度は落ちていません。"),
    (PALE, "通知が数分遅れることがある",
     "GitHub側の混み具合により、決まった時刻から数分ずれることがあります。仕様上避けられません。"),
    (PALE, "気象庁側の仕様変更で止まる可能性",
     "公式APIではないため、突然データが取れなくなり得ます。その場合もOpen-Meteoだけで動き続けます。"),
]):
    y = 1.72 + i * 1.03
    ellipse(s, 0.68, y + 0.16, 0.2, 0.2, fill=dotc)
    text(s, title, 1.12, y + 0.02, 4.6, 0.36, size=14.5, bold=True, color=WHITE)
    text(s, desc, 5.9, y + 0.02, 6.8, 0.8, size=11.5, color=PALE, spacing=18)

out = sys.argv[1] if len(sys.argv) > 1 else "rain-alert-overview.pptx"
prs.save(out)
print("生成完了:", out, f"({len(prs.slides._sldIdLst)} スライド)")
